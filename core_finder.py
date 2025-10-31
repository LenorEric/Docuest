import os
import re
import json
import time
import yaml
import hashlib
import sqlite3
import numpy as np
import requests
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional, Tuple

# Optional extractors
try:
    import textract
except Exception:
    textract = None
try:
    from docx import Document
except Exception:
    Document = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None
try:
    from pdfminer.high_level import extract_text as pdf_extract_text
except Exception:
    pdf_extract_text = None
try:
    import pandas as pd
except Exception:
    pd = None


# ---------------- Utils ----------------
def sha256_file(path: Path, bufsize: int = 1024 * 1024) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(bufsize), b""):
            h.update(chunk)
    return h.hexdigest()

def chunk_text(text: str, max_chars: int, overlap: int) -> List[Tuple[int, str]]:
    # normalize whitespace
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    if len(text) <= max_chars:
        return [(0, text)]
    out = []
    i, idx = 0, 0
    step = max_chars - overlap
    while i < len(text):
        part = text[i:i + max_chars]
        out.append((idx, part))
        idx += 1
        i += step
    return out

def cosine(a: np.ndarray, b: np.ndarray) -> float:
    denom = float(np.linalg.norm(a) * np.linalg.norm(b))
    if denom == 0.0:
        return 0.0
    return float(np.dot(a, b) / denom)

def parent_path_context(p: Path, depth: int = 3, joiner: str = " / ") -> str:
    parts = []
    cur = p.parent
    for _ in range(depth):
        if not cur or cur.name == "":
            break
        parts.append(cur.name)
        cur = cur.parent
    return joiner.join(reversed(parts))


# ---------------- Text extractors ----------------
def extract_text_any(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext == ".txt":
            try:
                return path.read_text(encoding="utf-8", errors="ignore")
            except Exception:
                return path.read_text(encoding="cp936", errors="ignore")
        if ext == ".pdf" and pdf_extract_text:
            # ignore extractable flag warnings explicitly
            return pdf_extract_text(str(path), check_extractable=False) or ""
        if ext == ".docx" and Document:
            doc = Document(str(path))
            paras = [p.text for p in doc.paragraphs]
            return "\n".join(paras)
        if ext == ".pptx" and Presentation:
            prs = Presentation(str(path))
            slides = []
            for i, s in enumerate(prs.slides, 1):
                t = []
                for shp in s.shapes:
                    if hasattr(shp, "text"):
                        t.append(shp.text)
                slides.append(f"[Slide {i}] " + " ".join(t))
            return "\n".join(slides)
        if ext in {".xlsx", ".xls"} and pd is not None:
            dfs = pd.read_excel(str(path), sheet_name=None, dtype=str)
            parts = []
            for sheet, df in dfs.items():
                df = df.fillna("")
                parts.append(f"[{sheet}]\n" + "\n".join(["\t".join(map(str, r)) for r in df.values]))
            return "\n".join(parts)
        if ext in {".doc", ".ppt", ".wps"} and textract is not None:
            return textract.process(str(path)).decode("utf-8", errors="ignore")
    except Exception:
        return ""
    return ""


# ---------------- SiliconFlow Embedding ----------------
def get_api_key(cfg: dict) -> str:
    path = cfg["embedding"].get("api_key_path", "./API_KEY.txt")
    if not os.path.exists(path):
        raise FileNotFoundError(f"API key file not found: {path}")
    return Path(path).read_text(encoding="utf-8").strip()

class SiliconFlowEmbedder:
    def __init__(self, cfg: dict):
        self.base = cfg["embedding"]["api_base"].rstrip("/")
        self.model = cfg["embedding"]["model"]
        self.api_key = get_api_key(cfg)
        self.timeout = int(cfg["embedding"].get("timeout_s", 60))
        self.retries = int(cfg["embedding"].get("max_retries", 5))

    def _post(self, data: dict) -> dict:
        url = f"{self.base}/embeddings"
        headers = {"Authorization": f"Bearer {self.api_key}", "Content-Type": "application/json"}
        for attempt in range(1, self.retries + 1):
            try:
                r = requests.post(url, json=data, headers=headers, timeout=self.timeout)
                if r.status_code == 200:
                    return r.json()
                print(f"[WARN] Embeddings HTTP {r.status_code}: {r.text[:200]}")
                if r.status_code in (429, 503):
                    time.sleep(2 ** attempt)
                    continue
                r.raise_for_status()
            except Exception as e:
                print(f"[ERROR] Embeddings attempt {attempt}: {e}")
                if attempt < self.retries:
                    time.sleep(2 ** attempt)
                else:
                    raise
        raise RuntimeError("SiliconFlow embeddings failed after retries")

    def embed_one(self, text: str) -> np.ndarray:
        data = self._post({"model": self.model, "input": text})
        vec = data["data"][0]["embedding"]
        return np.asarray(vec, dtype=np.float32)


# ---------------- SQLite schema ----------------
SCHEMA = """
CREATE TABLE IF NOT EXISTS documents(
  path TEXT PRIMARY KEY,
  file_hash TEXT,
  model TEXT,
  updated REAL
);
CREATE TABLE IF NOT EXISTS chunks(
  id INTEGER PRIMARY KEY AUTOINCREMENT,
  path TEXT,
  type TEXT DEFAULT 'content',   -- 'content' | 'filename' | 'path'
  text TEXT,
  embedding BLOB
);
"""

def open_db(sqlite_path: str) -> sqlite3.Connection:
    conn = sqlite3.connect(sqlite_path)
    conn.executescript(SCHEMA)
    # migrate: ensure 'type' column exists
    cols = {row[1] for row in conn.execute("PRAGMA table_info(chunks)")}
    if "type" not in cols:
        conn.execute("ALTER TABLE chunks ADD COLUMN type TEXT DEFAULT 'content'")
    return conn


# ---------------- Indexing ----------------
def index_all(cfg_path: str):
    cfg = yaml.safe_load(open(cfg_path, "r", encoding="utf-8"))
    embedder = SiliconFlowEmbedder(cfg)
    conn = open_db(cfg["storage"]["sqlite_path"])

    exts = {e.lower() for e in cfg["allowed_exts"]}
    dirs = cfg["include_dirs"]
    path_depth = int(cfg.get("path_context", {}).get("depth", 3))
    path_joiner = str(cfg.get("path_context", {}).get("joiner", " / "))

    files = []
    for d in dirs:
        for root, _, fs in os.walk(d):
            for f in fs:
                p = Path(root) / f
                if p.suffix.lower() in exts:
                    files.append(p)
    print(f"Found {len(files)} files.")

    model = embedder.model
    for f in files:
        try:
            h = sha256_file(f)
            # global dedup: same content + same model => skip
            cur = conn.execute(
                "SELECT path FROM documents WHERE file_hash=? AND model=?",
                (h, model)
            ).fetchone()
            if cur:
                print(f"[SKIP][DUPLICATE] {f} == {cur[0]}")
                continue

            text = extract_text_any(f)  # may be empty

            # Prepare chunks: content chunks if any text, plus filename and path context always
            to_insert: List[Tuple[str, str, bytes]] = []  # (type, preview_text, emb_blob)

            # content chunks
            if text:
                for _, ch in chunk_text(text, 1800, 200):
                    emb = embedder.embed_one(ch)
                    to_insert.append(("content", ch[:400], emb.tobytes()))
                indexed_via_content = True
            else:
                indexed_via_content = False

            # filename chunk
            fname = f.name
            emb_fname = embedder.embed_one(fname)
            to_insert.append(("filename", fname, emb_fname.tobytes()))

            # path-context chunk
            ctx = parent_path_context(f, depth=path_depth, joiner=path_joiner)
            if ctx:
                emb_ctx = embedder.embed_one(ctx)
                to_insert.append(("path", ctx, emb_ctx.tobytes()))

            # write DB
            conn.execute("DELETE FROM chunks WHERE path=?", (str(f),))
            for typ, prev, blob in to_insert:
                conn.execute(
                    "INSERT INTO chunks(path,type,text,embedding) VALUES(?,?,?,?)",
                    (str(f), typ, prev, blob)
                )
            conn.execute(
                "REPLACE INTO documents(path,file_hash,model,updated) VALUES(?,?,?,?)",
                (str(f), h, model, time.time())
            )
            conn.commit()
            tag = "" if indexed_via_content else " [name-only]"
            print(f"Indexed {f}{tag}")
        except Exception as e:
            print("Fail", f, e)
    conn.close()


# ---------------- DeepSeek chat ----------------
def call_deepseek(cfg: dict, query: str, context_text: str, can_render_markdown: bool) -> Tuple[str, Dict]:
    api_key = get_api_key(cfg)
    url = "https://api.siliconflow.cn/v1/chat/completions"

    # If GUI cannot render Markdown, force plain text.
    preface = "" if can_render_markdown else "仅回复纯文本。不要使用任何Markdown标记。"
    payload = {
        "model": "deepseek-ai/DeepSeek-V3.2-Exp",
        "messages": [
            {
                "role": "user",
                "content": f"{preface}\n请阅读以下材料来解释（几乎只以材料内容为依据，不添加额外信息，如果材料内容不充足，则回复较少的信息）：”{query}“\n\n材料：\n{context_text}"
            }
        ]
    }
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    r = requests.post(url, json=payload, headers=headers, timeout=180)
    try:
        data = r.json()
    except Exception:
        print(f"[DeepSeek HTTP {r.status_code}] {r.text[:200]}")
        raise

    # Print token usage to console
    usage = data.get("usage", {}) or {}
    print(f"[DeepSeek usage] prompt_tokens={usage.get('prompt_tokens')} "
          f"completion_tokens={usage.get('completion_tokens')} "
          f"total_tokens={usage.get('total_tokens')}")

    if r.status_code != 200:
        msg = data
        return f"[DeepSeek error {r.status_code}] {msg}", usage

    content = (data.get("choices", [{}])[0]
                  .get("message", {})
                  .get("content", "")) or ""
    # reasoning_content is optional; not displayed by default
    return content, usage


# ---------------- Search + Context ----------------
def search_docs(cfg_path: str, query: str, can_render_markdown: bool) -> Dict:
    cfg = yaml.safe_load(open(cfg_path, "r", encoding="utf-8"))
    embedder = SiliconFlowEmbedder(cfg)
    conn = open_db(cfg["storage"]["sqlite_path"])

    qvec = embedder.embed_one(query)

    # load all chunks
    rows = conn.execute("SELECT path, type, text, embedding FROM chunks").fetchall()
    if not rows:
        conn.close()
        return {"answer_md": "No index data.", "matched": []}

    scored: List[Tuple[float, str, str, str]] = []  # (score, path, type, text)
    for path, typ, txt, blob in rows:
        v = np.frombuffer(blob, dtype=np.float32)
        s = cosine(qvec, v)
        scored.append((s, path, typ, txt))

    conn.close()
    scored.sort(key=lambda x: x[0], reverse=True)

    # Filter ≥ 0.4 for context pool
    pool = [r for r in scored if r[0] >= 0.4]

    # Pick up to 10 best chunks overall; cap context text ~8k chars
    ctx_parts = []
    used = 0
    for s, p, typ, txt in pool[:50]:  # soft cap before length cut
        header = f"[{Path(p).name} | {typ} | {s:.3f}]"
        block = f"{header}\n{txt}"
        ctx_parts.append(block)
        used += len(block)
        if used >= 8000:
            break
        if len(ctx_parts) >= 10:
            break
    if not ctx_parts:
        ctx_parts = [f"[No chunk ≥ 0.4. Top1]\n{scored[0][3]}"]

    context_text = "\n\n".join(ctx_parts)

    # Ask DeepSeek
    answer_md, usage = call_deepseek(cfg, query, context_text, can_render_markdown)

    # Build matched file list (>0.3)
    matched = []
    seen = set()
    # also flag files that only had filename/path chunks (no content embeddings)
    # Strategy: if no 'content' chunk among top chunks for that path, mark as [filename-only?]
    # We approximate using entire scored list threshold >0.3.
    path_has_content: Dict[str, bool] = {}

    # Precompute which paths have any content chunk at all in DB (regardless of similarity)
    # Quick pass over rows we already have in memory: check type == 'content'
    # We didn't retain that, so we infer from 'scored' by scanning all entries
    for _, p, typ, _ in scored:
        if typ == "content":
            path_has_content[p] = True

    for s, p, typ, _ in scored:
        if s <= 0.3:
            break  # scored is sorted desc
        if p in seen:
            continue
        label = Path(p).name
        tag = "" if path_has_content.get(p, False) else " [filename-only]"
        matched.append(f"{s:.3f}  {p}{tag}  ({label})")
        seen.add(p)
        if len(matched) >= 100:
            break

    return {
        "answer_md": answer_md,
        "matched": matched,
        "usage": usage
    }
